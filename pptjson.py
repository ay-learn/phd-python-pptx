#!/usr/bin/env python3
import json
import os
import random
import string
import sys

import cairosvg
import cv2
import pptx
import yaml
from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# add
from pptx.enum.lang import MSO_LANGUAGE_ID
from lxml import etree
from pptx.oxml.xmlchemy import OxmlElement

SHRINK_NONE = MSO_AUTO_SIZE.NONE
SHRINK_TEXT = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
SHRINK_SHAPE = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

ALIGN_V_TOP = MSO_VERTICAL_ANCHOR.TOP
ALIGN_V_MIDDLE = MSO_VERTICAL_ANCHOR.MIDDLE
ALIGN_V_BOTTOM = MSO_VERTICAL_ANCHOR.BOTTOM
ALIGN_V_MIXED = MSO_VERTICAL_ANCHOR.MIXED

MSO_VERTICAL_ANCHOR.MIDDLE

ALIGN_H_CENTER = PP_ALIGN.CENTER
ALIGN_H_LEFT = PP_ALIGN.LEFT
ALIGN_H_RIGHT = PP_ALIGN.RIGHT
ALIGN_H_JUSTIFY = PP_ALIGN.JUSTIFY_LOW
# ALIGN_H_JUSTIFY = PP_ALIGN.JUSTIFY
ALIGN_H_DISTRIBUTE = PP_ALIGN.DISTRIBUTE
ALIGN_H_THAI_DISTRIBUTE = PP_ALIGN.THAI_DISTRIBUTE
ALIGN_H_MIXED = PP_ALIGN.MIXED

BLUE_DARK1 = RGBColor(30, 4, 91)
BLUE_DARK2 = RGBColor(0, 32, 96)
RED = RGBColor(255, 0, 0)
GREEN = RGBColor(0, 135, 0)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)

IMG_BUTTOM = "IMG_BUTTOM"
IMG_BESIDE = "IMG_RIGHT"


def generate_random_string(length=6):
    letters = string.ascii_lowercase
    return "".join(random.choice(letters) for _ in range(length))


def get_image_dpi(image_path):
    try:
        with Image.open(image_path) as img:
            dpi = img.info.get("dpi")
            print("image_path", image_path, "DPI: ", dpi)
            if dpi:
                return dpi
            else:
                return None
    except Exception:
        return None


def get_img_width(image_path, image_position):
    if image_path is None or image_path == "":
        print(f"image_path={image_path}", file=sys.stderr)
        return Inches(0), image_path
    if not os.path.isfile(image_path):
        print(f"File {image_path} does not found", file=sys.stderr)
        return Inches(0), image_path

    if image_position is None:
        image_position == IMG_BUTTOM

    if image_path.lower().endswith(".svg"):
        rand_name = generate_random_string(6)
        output_path = f"./build/{rand_name}.png"
        cairosvg.svg2png(url=image_path, write_to=output_path)
        image_path = output_path

    image_np = cv2.imread(image_path)

    image_width, image_height = 0, 0
    try:
        image_height, image_width, _ = image_np.shape
    except AttributeError as e:
        print("image_path=", image_path, "\n", e)
        return Inches(0), image_path

    # PIL way
    # image = Image.open(image_path)
    # image_width, image_height = image.size

    xy = get_image_dpi(image_path)
    if xy is not None:
        dpi_x, dpi_y = xy[0], xy[1]
    else:
        dpi_x, dpi_y = 100, 100
    image_width_inches = image_width / dpi_x
    image_height_inches = image_height / dpi_y

    width = Inches(4) * image_width_inches / image_height_inches
    if image_position == IMG_BESIDE:
        if width > Inches(5):
            width = Inches(5)
    else:
        if width > Inches(14):
            width = Inches(14)

    return width, image_path


# Define a base class for text elements
class TextElement:
    def __init__(self, shape):
        self.shape = shape

    def text(self, text):
        self.shape.text = text
        return self

    def width(self, inches):
        self.shape.width = inches
        return self

    def height(self, inches):
        self.shape.height = inches
        return self

    def X(self, inches):
        self.shape.top = inches
        return self

    def Y(self, inches):
        self.shape.left = inches
        return self

    def color(self, rgb):
        for paragraph in self.shape.text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                font.color.rgb = rgb
        return self

    def upper(self):
        self.shape.text = self.shape.text.upper()
        return self

    def bold(self):
        for paragraph in self.shape.text_frame.paragraphs:
            paragraph.font.bold = True
        return self

    def font_size(self, points):
        for paragraph in self.shape.text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                font.size = Pt(points)
        return self

    def font_name(self, name):
        for paragraph in self.shape.text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                font.name = name
                font.language_id = LANG  # MSO_LANGUAGE_ID.ARABIC
        return self

    def align_H(self, alignment):
        for paragraph in self.shape.text_frame.paragraphs:
            paragraph.alignment = alignment
            paragraph.level = LEVEL
        return self

    def align_V(self, anchor):
        self.shape.text_frame.vertical_anchor = anchor
        return self

    def shrink(self, auto_size):
        for paragraph in self.shape.text_frame.paragraphs:
            paragraph.auto_size = auto_size
        return self


# Define a subclass for title element
class Subtitle(TextElement):
    def __init__(self, shape):
        super().__init__(shape)


class Title(TextElement):
    def __init__(self, shape):
        super().__init__(shape)


class TextBox(TextElement):
    def __init__(self, shape):
        super().__init__(shape)


class Paragraph(TextElement):
    def __init__(self, shape):
        super().__init__(shape)

    def add_paragraph(self, text="", color=None, font_size=None, font_name=None):
        paragraph = self.shape.text_frame.add_paragraph()

        # paragraph.alignment = ALIGN_H_JUSTIFY
        paragraph.level = LEVEL
        paragraph.alignment = ALIGN

        run = paragraph.add_run()
        run.text = text
        run.level = LEVEL

        if color:
            run.font.color.rgb = color
        if font_size:
            run.font.size = Pt(font_size)
        if font_name:
            run.font.name = font_name
        return TextElement(paragraph)


class Slide:
    def __init__(self, presentation):
        self.presentation = presentation

        self.presentation.slide_width = Inches(16)
        self.presentation.slide_height = Inches(9)

        self.slide = presentation.slides.add_slide(presentation.slide_layouts[0])

        # Create attributes for each element as instances of their respective classes
        self.title_shape = Title(self.slide.shapes.title)
        self.subtitle_shape = Subtitle(self.slide.placeholders[1])
        self.subtitle_shape_2 = Subtitle(self.slide.placeholders[5])
        self.page_shape_3 = Subtitle(self.slide.placeholders[6])
        self.paragraph_shape = Paragraph(self.slide.placeholders[2])
        self.paragraph_shape_2 = Paragraph(self.slide.placeholders[3])
        self.paragraph_shape_3 = Paragraph(self.slide.placeholders[4])

        # Create an attribute for the shapes collection
        self.shapes = self.slide.shapes

    # Define methods to access the attributes of the elements
    def add_title(self, text=None):  # Rename this method
        if text is not None:
            self.title_shape.text(text)

        # Return the title_shape attribute to use its methods with dot notation
        return self.title_shape

    def add_subtitle(self, text=None):
        if text is not None:
            self.subtitle_shape.text(text)
        return self.subtitle_shape

    def add_page(self, text=None):
        if text is not None:
            self.page_shape_3.text(text)
        return self.page_shape_3

    def add_paragraph(self, text: str = "", color=None, font_size=None, font_name=None):
        return self.paragraph_shape.add_paragraph(text, color, font_size, font_name)

    def set_paragraph(self, width=None, height=None, X=None, Y=None, rgb=None):
        if width is not None:
            self.paragraph_shape.width(width)
        if height is not None:
            self.paragraph_shape.height(height)
        if X is not None:
            self.paragraph_shape.X(Y)
        if Y is not None:
            self.paragraph_shape.Y(Y)
        if rgb is not None:
            self.paragraph_shape.color(rgb)
        return self.paragraph_shape

    def add_image(self, image_path=None, image_position=None):
        if image_path is None or image_path == "":
            print(f"image_path={image_path}", file=sys.stderr)
            return self
        if not os.path.isfile(image_path):
            print(f"File {image_path} does not found", file=sys.stderr)
            return self

        height = Inches(4)
        width, image_path = get_img_width(image_path, image_position)

        if image_position == IMG_BESIDE:
            if ARABIC:
                left = Inches(0.25)
                top = Inches(3)
            else:
                left = Inches(16) - width - Inches(0.25)
                top = Inches(3)

            self.shapes.add_picture(image_path, left=left, top=top, width=width)

        else:
            left = (Inches(16) - width) / 2
            top = Inches(4.30)

            self.shapes.add_picture(
                image_path, left=left, top=top, width=width, height=height
            )
        return self

    def save(self, path):
        if self.presentation is not None:
            self.presentation.save(path)
            print(f"---------\nsave => {path}")
        else:
            exit("Presentation not saved")


def Presentation(path):
    return pptx.Presentation(path)


def remove_first_slide(presentation):
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[0])


def load_slides_from_json(json_file_path):
    with open(json_file_path, "r") as f:
        slides_data = json.load(f)

    for slide_data in slides_data:
        add_slide_from_data(slide_data)


def load_slides_from_yaml(yaml_file_path):
    with open(yaml_file_path, "r") as f:
        slides_data = yaml.safe_load(f)

    for slide_data in slides_data:
        add_slide_from_data(slide_data)


def add_slide_from_data(slide_data):
    slide = Slide(presentation)

    title = slide_data.get("title", "")
    subtitle = slide_data.get("subtitle", "")
    subtitle_2 = slide_data.get("subtitle_2", "")
    slide_data.get("page", "")
    paragraphs = slide_data.get("paragraphs", [])
    paragraphs_2 = slide_data.get("paragraphs_2", [])  # noqa: F841
    paragraphs_3 = slide_data.get("paragraphs_3", [])  # noqa: F841
    image = slide_data.get("image_path", "")
    image_position = slide_data.get("image_position", IMG_BUTTOM)

    img_width, _ = get_img_width(image, image_position)

    if image_position == IMG_BESIDE:
        par_width = int(Inches(15) - img_width - Inches(0.50))

        if par_width < 0:
            print("img_width is bigger")
            par_width = 0
        print(
            f"IMG_BESIDE: 15 Inches({Inches(15)}) - img_width({img_width}) \tpar_width: {par_width}"
        )

        if ARABIC:
            TEXT_START = int(img_width + Inches(0.5))
        else:
            TEXT_START = Inches(1)
    else:
        par_width = Inches(14.5)
        TEXT_START = Inches(1)

    if title:
        (
            slide.add_title(title)
            .X(Inches(0.5))
            .Y(Inches(0.5))
            .width(Inches(15))
            .height(Inches(1))
            .upper()
            .color(RED)
            .bold()
            .font_size(48)
            .font_name("Arial")
            .align_H(ALIGN_H_CENTER)
            .align_V(ALIGN_V_TOP)
            .shrink(SHRINK_SHAPE)
        )

    if subtitle:
        (
            slide.add_subtitle(subtitle)
            .X(Inches(2))
            .Y(Inches(1))
            .width(Inches(14))
            .height(Inches(1))
            .bold()
            .color(GREEN)
            .font_size(36)
            .font_name("Monotype Corsiva")
            .align_H(ALIGN)
            .align_V(ALIGN_V_TOP)
            .shrink(SHRINK_TEXT)
        )

    if paragraphs:
        (
            slide.set_paragraph()
            .width(par_width)
            .height(Inches(6))
            .X(Inches(2.5))
            .Y(TEXT_START)
            .font_size(PARAGAPH_FONT_SIZE)
            .font_name("Arial")
            .align_H(ALIGN)
        )
        for paragraph in paragraphs:
            slide.add_paragraph(
                text=paragraph,
                color=BLACK,
                font_size=PARAGAPH_FONT_SIZE,
                font_name="Arial",
            )

    if image:
        slide.add_image(image_path=image, image_position=image_position)

    return slide


def RtoL(shape):
    textb = shape.text_frame._txBody
    for bad in textb.xpath("./a:p/a:pPr"):
        bad.set("rtl", "1")


def set_rtl(presentation):
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                textb = shape.text_frame._txBody
                for bad in textb.xpath("./a:p/a:pPr"):
                    bad.set("rtl", "1")


def del_empty_palaceholder(presentation):
    for slide in presentation.slides:
        for placeholder in slide.placeholders:
            if not placeholder.text_frame.text:  # check if the text is empty
                placeholder.element.delete()


ARABIC = 1  # True/FALSE

if ARABIC:
    print("ARABIC")
    LANG = MSO_LANGUAGE_ID.ARABIC
    ALIGN = ALIGN_H_RIGHT
else:
    print("ENGLISH")
    LANG = MSO_LANGUAGE_ID.ENGLISH_US
    ALIGN = ALIGN_H_LEFT
LEVEL = 0

PARAGAPH_FONT_SIZE = 26

presentation = Presentation("t2.pptx")

# load_slides_from_json("ppt1.json")
load_slides_from_yaml("f4.yaml")

remove_first_slide(presentation)
if ARABIC:
    set_rtl(presentation)
del_empty_palaceholder(presentation)
presentation.save("new_slide.pptx")
