import collections.abc  # noqa: F401

import pptx
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches
from pptx.util import Pt

FIT_NONE = MSO_AUTO_SIZE.NONE

ALIGN_TOP = MSO_ANCHOR.TOP
ALIGN_CENTER = PP_ALIGN.CENTER
ALIGN_LEFT = PP_ALIGN.LEFT

BLUE_DARK1 = RGBColor(30, 4, 91)
BLUE_DARK2 = RGBColor(0, 32, 96)
RED = RGBColor(255, 0, 0)
GREEN = RGBColor(0, 135, 0)
BLACK = RGBColor(0, 0, 0)

IMG_RIGHT = True
IMG_BUTTOM = None


def print_placeholder(presentation):
    for i in range(99):
        print(f"--------{i}--------")
        master_layout = i
        try:
            slide = presentation.slides.add_slide(
                presentation.slide_layouts[master_layout]
            )
            # for shape in slide.placeholders:
            #     print("%d %s" % (shape.placeholder_format.idx, shape.name))
            for shape in slide.shapes:
                # print('%s' % shape.shape_type)
                if shape.is_placeholder:
                    phf = shape.placeholder_format
                    print("%d, %s" % (phf.idx, phf.type))
            print("_________________")
        except Exception:
            break


def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element


def makeParaBulletPointed(para, symbol="●"):
    """Bullets are set to Arial,
    actual text can be a different font"""
    pPr = para._p.get_or_add_pPr()
    ## Set marL and indent attributes
    pPr.set("marL", "171450")
    # pPr.set("indent", "171450")
    pPr.set("indent", "0")
    ## Add buFont
    _ = SubElement(
        parent=pPr,
        tagname="a:buFont",
        typeface="Arial",
        panose="020B0604020202020204",
        pitchFamily="34",
        charset="0",
    )
    ## Add buChar
    _ = SubElement(parent=pPr, tagname="a:buChar", char=symbol)


def set_options(
    shape,
    text=None,
    upper=None,
    width=None,
    height=None,
    X=None,
    Y=None,
    color=None,
    bold=None,
    font_size=None,
    font_name=None,
    align_H=None,
    align_V=None,
    shrink=None,
) -> None:
    if text is not None:
        shape.text = text
        if upper is not None:
            shape.text = text.upper()
    if width is not None:
        shape.width = width
    if height is not None:
        shape.height = height
    if X is not None:
        shape.top = X
    if Y is not None:
        shape.left = Y
    for paragraph in shape.text_frame.paragraphs:
        if color is not None:
            paragraph.font.color.rgb = color
        if bold is not None:
            paragraph.font.bold = bold
        if font_size is not None:
            paragraph.font.size = Pt(font_size)
        if font_name is not None:
            paragraph.font.name = font_name
        if align_H is not None:
            paragraph.alignment = align_H
    if align_V is not None:
        shape.text_frame.vertical_anchor = align_V

    if shrink is not None:
        shape.text_frame.auto_size = shrink


def add_paragraph_bac(text_box=None, text=None, symbol="●") -> None:
    if text_box is None or text is None:
        exit("missing args in add_paragraph")
    paragraph = text_box.font.size
    paragraph = text_box.add_paragraph()
    paragraph.text = f" {text}"
    makeParaBulletPointed(paragraph, symbol)


def add_paragraph(text_box=None, text=None, symbol="●") -> None:
    if text_box is None or text is None:
        exit("missing args in add_paragraph")
    paragraph = text_box.add_paragraph()
    paragraph.text = text


def add_slide(
    presentation,
    title=None,
    subtitle=None,
    symbol="●",
    texts=None,
    image_path=None,
    image_position=None,
) -> None:
    presentation.slide_width = Inches(16)
    presentation.slide_height = Inches(9)

    # print_placeholder(presentation)
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])

    ##1 text_box = slide.placeholders[2]
    ##1 text_frame = text_box.text_frame
    ##1
    ##1 paragraph1 = text_frame.add_paragraph()
    ##1 paragraph1.text = "AAA"
    ##1 paragraph2 = text_frame.add_paragraph()
    ##1 paragraph2.text = "BBB"
    ##1 print("fin")
    # return presentation

    # TODO: if placeholder = ... do ... elif do ...elif do...

    if title is not None:
        title_shape = slide.shapes.title

        set_options(
            title_shape,
            text=title,
            upper=True,
            width=Inches(16),
            height=None,
            X=Inches(0.5),
            Y=Inches(0),
            color=RED,
            bold=True,
            font_size=36,
            font_name="Arial",
            align_H=ALIGN_CENTER,
            align_V=ALIGN_TOP,
            shrink=FIT_NONE,
        )

    if subtitle is not None:
        subtitle_shape = slide.placeholders[1]
        set_options(
            subtitle_shape,
            text=subtitle,
            upper=None,
            width=Inches(8),
            height=None,
            X=Inches(2),
            Y=Inches(1),
            color=GREEN,
            bold=True,
            font_size=36,
            font_name="Monotype Corsiva",
            align_H=ALIGN_LEFT,
            align_V=ALIGN_TOP,
            shrink=FIT_NONE,
        )
    if texts is not None:
        text_box = slide.placeholders[2]
        text_frame = text_box.text_frame
        for text in texts:
            add_paragraph(text_frame, symbol=symbol, text=text)

        if image_position is IMG_RIGHT:
            img_width = Inches(4)
        else:
            img_width = Inches(0)

        set_options(
            text_box,
            text=None,
            upper=None,
            width=Inches(15) - img_width,
            height=None,
            X=Inches(2.50),
            Y=Inches(1),
            color=BLACK,
            bold=False,
            font_size=24,
            font_name="Arial",
            align_H=ALIGN_LEFT,
            align_V=ALIGN_TOP,
            shrink=FIT_NONE,
        )

    if image_path is not None:
        shapes = slide.shapes
        if image_position is IMG_RIGHT:
            left = Inches(12)
            top = Inches(3)
            width = Inches(3.90)
        else:
            left = Inches(6.5)
            top = Inches(5.30)
            width = None
        height = Inches(3)

        shapes.add_picture("img.png", left, top, width=width, height=height)

    return presentation


# 1,3,7 and ... has bullet
if __name__ == "__main__":
    # presentation = pptx.Presentation("template1.pptx")
    presentation = pptx.Presentation("/tmp/t8.pptx")
    # presentation = pptx.Presentation("candy.pptx")
    # presentation = pptx.Presentation()
    add_slide(
        presentation,
        title="Introduction",
        subtitle="Définitions",
        symbol="-",
        texts=[
            "Une source radioactive est une quantité connue d'un radionucléide qui émet un rayonnement ionisant.",
            "Un rayonnement ionisant est un rayonnement électromagnétique ou corpusculaire capable de produire directement ou indirectement des ions lors de son passage à travers la matière. ",
            "CCC",
        ],
        image_path="img.png",
        image_position=IMG_BUTTOM,  # IMG_RIGHT
    )


    add_slide(
        presentation,
        title="Introduction1",
        subtitle="Définitions1",
        symbol="-",
        texts=[
            "Une source radioactive est une quantité connue d'un radionucléide qui émet un rayonnement ionisant.",
            "Un rayonnement ionisant est un rayonnement électromagnétique ou corpusculaire capable de produire directement ou indirectement des ions lors de son passage à travers la matière. ",
        ],
        image_path="img.png",
        image_position=IMG_BUTTOM,  # IMG_RIGHT
    )

    if presentation is not None:
        presentation.save("new_slide.pptx")
    else:
        exit("Presentation not saved")
